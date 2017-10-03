#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Main script for udemy-dl."""
from __future__ import unicode_literals
from __future__ import print_function
from builtins import str
from builtins import input

import json
import re
import os
import sys
import getpass
import argparse
import errno
import time
import shutil
import logging
from io import open, BytesIO
from distutils.spawn import find_executable

try:
    from urllib import unquote
except ImportError:
    from urllib.parse import unquote

import requests
import requests.sessions
import colorlog
import pptx
from slugify import slugify, SLUG_OK

from .download import download, DLException, USER_AGENT
from ._version import __version__


# global variable
debug = False
debug_path = ''
use_ffmpeg = False
downloader = 'aria2c'
video_quality = '654321'
quality_list = ['1080', '720', '640', '480', '360', '240']
logger = colorlog.getLogger(__name__)

# udemy.com api
LOGOUT_URL = 'http://www.udemy.com/user/logout'
LOGIN_POPUP_URL = 'https://www.udemy.com/join/login-popup'
CHECK_COURSE_STATUS_URL = 'https://www.udemy.com/api-2.0/users/me/course-previews/?course={course_id}'
QUIZ_URL = 'https://www.udemy.com/api-2.0/quizzes/{quiz_id}/assessments?fields[assessment]=@all&page_size=250'
ATTACHED_FILE_URL = 'https://www.udemy.com/api-2.0/users/me/subscribed-courses/{course_id}/lectures/{lecture_id}/supplementary-assets/{asset_id}?fields[asset]=download_urls'
LOGIN_URL = 'https://www.udemy.com/join/login-popup/?displayType=ajax&display_type=popup&showSkipButton=1&returnUrlAfterLogin=https%3A%2F%2Fwww.udemy.com%2F&next=https%3A%2F%2Fwww.udemy.com%2F&locale=en_US'
GET_LECTURE_URL = 'https://www.udemy.com/api-2.0/users/me/subscribed-courses/{course_id}/lectures/{lecture_id}?fields[asset]=@min,download_urls,external_url,slide_urls&fields[course]=id,is_paid,url&fields[lecture]=@default,view_html,course&page_config=ct_v4'
COURSE_INFO_URL = 'https://www.udemy.com/api-2.0/courses/{course_id}/cached-subscriber-curriculum-items?fields[asset]=@min,title,filename,asset_type,external_url,length&fields[chapter]=@min,description,object_index,title,sort_order&fields[lecture]=@min,object_index,asset,supplementary_assets,sort_order,is_published,is_free&fields[quiz]=@min,object_index,title,sort_order,is_published&page_size=550'
COURSE_TITLE_URL = 'https://www.udemy.com/api-2.0/courses/{course_id}?fields[course]=title'


def logging_exception(type_, value, tb):
    """Catch Exception message."""
    logger.error("Exception",
                 exc_info=(type_, value, tb))
    # sys.__excepthook__(type_, value, tb)


# Install exception handler
sys.excepthook = logging_exception


class Session:

    """Starting session with proper headers to access udemy site."""

    headers = {'User-Agent': USER_AGENT,
               'X-Requested-With': 'XMLHttpRequest',
               'Host': 'www.udemy.com',
               'Referer': LOGIN_POPUP_URL}

    def __init__(self):
        """Init session."""
        self.session = requests.sessions.Session()

    def set_auth_headers(self, access_token, client_id):
        """Setting up authentication headers."""
        self.headers['X-Udemy-Bearer-Token'] = access_token
        self.headers['X-Udemy-Client-Id'] = client_id
        self.headers['Authorization'] = "Bearer " + access_token
        self.headers['X-Udemy-Authorization'] = "Bearer " + access_token

    def get(self, url):
        """Retrieving content of a given url."""
        return self.session.get(url, headers=self.headers)

    def post(self, url, data):
        """HTTP post given data with requests object."""
        return self.session.post(url, data, headers=self.headers)


session = Session()


def safeencode(unsafetext):
    """safe encode filenames."""
    text = slugify(unsafetext, lower=False, spaces=True, ok=SLUG_OK + '().')
    return text


def save_debug_data(debug_data, debug_name, ext):
    """Save debug data to find bugs."""
    debug_str = str(debug_data)
    debug_str = re.sub(r"[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+", 'USER@DOMAIN.COM', debug_str)
    debug_time = time.strftime("%Y%m%d-%H%M%S")
    debug_file_name = os.path.join(debug_path, 'DEBUG-{0}-{1}.{2}'.format(debug_name, debug_time, ext))

    with open(debug_file_name, 'w', encoding='utf-8') as save_debug:
        save_debug.write(debug_str)


def get_csrf_token():
    """Extractig CSRF Token from login page."""
    try:
        response = session.get(LOGIN_POPUP_URL)
        match = re.search(r"name='csrfmiddlewaretoken'\s+value='(.*)'", response.text)
        return match.group(1)
    except AttributeError:
        session.get(LOGOUT_URL)
        response = session.get(LOGIN_POPUP_URL)
        match = re.search(r"name='csrfmiddlewaretoken'\s+value='(.*)'", response.text)
        return match.group(1)


def login(username, password):
    """Login with popup-page."""
    logger.info("Trying to log in ...")
    csrf_token = get_csrf_token()
    payload = {'isSubmitted': 1, 'email': username, 'password': password,
               'displayType': 'ajax', 'csrfmiddlewaretoken': csrf_token}
    response = session.post(LOGIN_URL, payload)

    access_token = response.cookies.get('access_token')
    client_id = response.cookies.get('client_id')
    response_text = response.text

    if '<li>You have exceeded the maximum number of requests per hour.</li>' in response_text:
        logger.error('You have exceeded the maximum number of login requests per hour.')
        sys.exit(1)

    elif '<li>Please check your email and password.</li>' in response_text:
        logger.error('Wrong Username or Password!')
        sys.exit(1)

    elif access_token is None:
        logger.error("Couldn't fetch token!")
        sys.exit(1)

    elif 'error' in response_text:
        print(response_text)
        logger.error('Found error in login page')
        sys.exit(1)

    session.set_auth_headers(access_token, client_id)

    logger.info("Login success.")


def get_course_id(course_link):
    """Retrieving course ID."""
    if 'udemy.com/draft/' in course_link:
        course_id = course_link.split('/')[-1]
        logger.info("Found draft...id: %s", course_id)
        return course_id

    response = session.get(course_link)
    response_text = response.text

    if 'data-purpose="take-this-course-button"' in response_text:
        logger.error('Please Enroll in this course')
        sys.exit(1)

    if debug:
        save_debug_data(response_text, 'get_course_id', 'html')

    logger.debug('Searching course id...')
    matches = re.search(r'data-course-id="(\d+)"', response_text, re.IGNORECASE)
    if matches:
        course_id = matches.groups()[0]
    else:
        matches = re.search(r'property="og:image"\s+content="([^"]+)"', response_text, re.IGNORECASE)
        course_id = matches.groups()[0].rsplit('/', 1)[-1].split('_', 1)[0] if matches else None

    if not course_id:
        logger.error('Course id not found!')
        sys.exit(1)
    else:
        logger.info('Found course id: %s', course_id)

    return course_id

def get_course_title(course_id):
    """Getting course title using course id"""
    course_title_url = COURSE_TITLE_URL.format(course_id=course_id)
    course_title_data = session.get(course_title_url).json()
    logger.debug('Saving course_title_data...')
    if debug:
        save_debug_data(course_title_data, 'get_course_title', 'txt')
    course_title = course_title_data['title']
    logger.info('Found course title: %s', course_title)
    return course_title

def extract_lecture_url(course_id, lecture_id, asset_type):
    """Extracting Lecture url."""
    get_url = GET_LECTURE_URL.format(course_id=course_id, lecture_id=lecture_id)
    json_source = session.get(get_url).json()

    logger.debug('Saving json source: %s', lecture_id)

    if debug:
        save_debug_data(json_source, 'json_source_' + str(lecture_id), 'txt')

    lecture_parser_dict = {'Video': parse_video,
                           'VideoMashup': parse_video,
                           'E-Book': parse_ebook,
                           'Presentation': parse_presentation,
                           'Article': (lambda _: (get_url, 'Article')),
                           'Audio': parse_audio,
                           'File': parse_file}
    parser = lecture_parser_dict.get(asset_type)

    if parser:
        return parser(json_source)
    else:
        logger.critical("Couldn't extract lecture url: %s", lecture_id)
        return (None, None)


def parse_video(json_source):
    """Extracting video URL from json_source for type Video and VideoMashup."""
    clean_html = json_source['view_html'].replace('&quot;', '"')

    # This regex should grab just the JSON in the videojs-setup-data attribute
    json_regex = 'videojs-setup-data\=\"(.*?)\"\\n\s+text-tracks\=\"'

    the_json = re.findall(json_regex, clean_html, re.DOTALL)[0]

    # Load the matched JSON
    try:
        data = json.loads(the_json)
        list_videos = [(s['src'], s['label']) for s in data['sources']]
    except ValueError:
        list_videos = []
        logger.debug('list_videos JSON could not be loaded')

    # TODO: These regexes probably won't work anymore since the HTML
    # has changed.
    list_m3u8 = re.findall(r'source\s+src="(.*?)".*?data-res="\w+[^\d]+"', json_source['view_html'], re.DOTALL)
    caption_link = re.search(r'(https:\/\/udemy-captions.*?\.vtt.+?)&#39', json_source['view_html'])
    caption_language = re.search(r'language&#39;: &#39;(.*?)&#39', json_source['view_html'])
    caption_list = [(unescape(caption_link.group(1)), caption_language.group(1))] if (caption_link and caption_language) else None
    dict_videos = {}
    found = False

    if json_source['asset']['download_urls']:
        logger.debug('Found Videos in json_source')
        for element in json_source['asset']['download_urls']['Video']:
            dict_videos[element['label']] = element['file']
        found = True

    if list_videos and not found:
        logger.debug('Found list_videos')
        for link, quality in list_videos:
            dict_videos[quality] = link
        found = True

    if list_m3u8 and not found:  # less priority to m3u8 for saving links
        logger.debug('Found list_m3u8')
        link = unescape(unquote(unquote(list_m3u8[0])))
        m3u8_str = requests.get(link).text

        for line in m3u8_str.splitlines():
            quality = re.match(r'.*/hls_(\d+)_', line)
            if quality:
                dict_videos.setdefault(quality.group(1), line)

    for item in quality_list:
        if dict_videos.get(item):
            return (dict_videos[item], 'Video', caption_list)
    logger.critical('Skipped. Expected quality not found!')
    return (None, None, None)


def parse_ebook(json_source):
    """Extracting URL from json_source for type E-book."""
    for element in json_source['asset']['download_urls']['E-Book']:
        if element['label'] == 'download':
            return (element['file'], 'Ebook')
        else:
            logger.critical("Skipped. Couldn't fetch e-book!")
            return (None, None)


def parse_audio(json_source):
    """Extracting URL from json_source for type Audio."""
    for element in json_source['asset']['download_urls']['Audio']:
        if element['label'] == 'download':
            return (element['file'], 'Audio')
        else:
            logger.critical("Skipped. Couldn't fetch audio!")
            return (None, None)


def parse_presentation(json_source):
    """Extracting URL from json_source for type Presentation."""
    old_pdf = json_source['asset'].get('download_urls', {})
    new_pptx = json_source['asset'].get('slide_urls', {})
    if old_pdf:
        for element in old_pdf.get('Presentation'):
            if element.get('label') == 'download':
                return (element['file'], 'Presentation')
    else:
        return (new_pptx, 'Pre_pptx')

    logger.critical("Skipped. Couldn't fetch presentation!")
    return (None, None)

def parse_file(json_source):
    """Extracting URL from json_source for type File."""
    for element in json_source['asset']['download_urls']['File']:
        if element['label'] == 'download':
            return (element['file'], 'File')
        else:
            logger.critical("Skipped. Couldn't fetch File!")
            return (None, None)


def get_quiz(link, filename):  # need better support
    """Get multiple choice."""
    logger.info('Getting Quiz')
    quiz = session.get(link).text
    htm = r'''<!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <title>quiz_title</title>
    </head>
    <body>

        <div style="background-color:#ffffaa" align = "center">
            <script>
                var c = '';
                var opt = '';
                var n = ['a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j']

                var quiz = json_data;

                for (var i = 0; i < quiz.results.length; i++) {

                    var question = '<font color="ff0000"><br><h3>Question-' + i + ': </h3> <b>' + quiz.results[i].prompt.question + '</b> <br>';

                    for (var j = 0; j < quiz.results[i].prompt.answers.length; j++) {

                        opt = opt.concat(quiz.results[i].prompt.answers[j] ? '<p style="cursor: pointer" onclick="this.style=\'font-weight: bold; color: #CE534D\'">' + n[j] + ': ' + quiz.results[i].prompt.answers[j].toString().replace(/,/g, '<br>') + '</p>' : '<br>');
                    }
                    var ans = quiz.results[i].correct_response ? '<p> <br> <input style="cursor: pointer" type="button" onclick="this.value=\'Answer is: ' + quiz.results[i].correct_response.toString() + '\'" value="Show Answer"></input> </p>' : '<br>';
                    c = c.concat(question, opt, ans);
                    opt = ''
                    }
                document.write(c);
            </script>
        </div>
    </body>
    </html>'''.replace('quiz_title', filename).replace('json_data', str(quiz))

    filename = str(filename) + '.html'

    dest_dir = os.path.join(".", "quizzes", filename)
    mkdir(os.path.dirname(dest_dir))

    with open(dest_dir, 'w', encoding='utf-8') as quiz_file:
        quiz_file.write(htm)


def parse_quiz(quiz_id):
    """Parse multiple choice."""
    quiz_url = QUIZ_URL.format(quiz_id=quiz_id)
    return [quiz_url, 'quiz']


def unescape(strs):
    """Replace HTML-safe sequences "&amp;", "&lt;"" and "&gt;" to special characters."""
    strs = strs.replace("&amp;", "&")
    strs = strs.replace("&lt;", "<")
    strs = strs.replace("&gt;", ">")
    return strs


def get_data_links(course_id, lecture_start, lecture_end):
    """Getting video links from api 2.0."""
    course_url = COURSE_INFO_URL.format(course_id=course_id)
    course_data = session.get(course_url).json()
    logger.debug('Saving course_data...')
    if debug:
        save_debug_data(course_data, 'get_course_data', 'txt')

    chapter = None
    data_list = []
    supported_asset_type = ['Video', 'E-Book', 'VideoMashup', 'Audio', 'Presentation', 'Article', 'File']

    lecture_number = 1
    chapter_number = 0
    item_count = 0.0
    course_data_len = len(course_data['results'])

    # A udemy course has chapters, each having one or more lectures
    for item in course_data['results']:
        item_count += 1
        lecture = item['title']
        lecture_id = item['id']

        update_progress(item_count / course_data_len, "Getting course data")

        if item['_class'] == 'chapter':
            chapter = item['title']
            chapter_number += 1

        elif item['_class'] == 'lecture' and item['asset']['asset_type'] in supported_asset_type:
            if valid_lecture(lecture_number, lecture_start, lecture_end):
                try:
                    if item['asset']['asset_type'] in ['Video', 'VideoMashup']:
                        data_url, data_type, caption_list = extract_lecture_url(course_id, lecture_id, item['asset']['asset_type'])
                    else:
                        caption_list = []
                        data_url, data_type = extract_lecture_url(course_id, lecture_id, item['asset']['asset_type'])

                    if data_url is None:
                        lecture_number += 1
                        continue

                    attached_list = []
                    if item.get('supplementary_assets'):
                        for assets in item['supplementary_assets']:
                            attached_list.append({'filename': assets['filename'],
                                                  'id': assets['id']})

                    attached_info = {'course_id': course_id,
                                     'lecture_id': lecture_id,
                                     'attached_list': attached_list}

                    data_list.append({'chapter': chapter,
                                      'lecture': lecture,
                                      'data_url': unescape(data_url),
                                      'data_type': data_type,
                                      'attached_info': attached_info,
                                      'caption_list': caption_list,
                                      'lecture_number': int(lecture_number),
                                      'chapter_number': int(chapter_number)})
                except Exception as exc:
                    logger.critical('Cannot download lecture "%s": "%s"', safeencode(lecture), exc)

            lecture_number += 1

        elif item['_class'] == 'quiz':
            data_url, data_type = parse_quiz(lecture_id)

            data_list.append({'chapter': chapter,
                              'lecture': lecture,
                              'data_url': data_url,
                              'data_type': data_type,
                              'attached_info': {'attached_list': []},
                              'caption_list': [],
                              'lecture_number': int(item['object_index']),
                              'chapter_number': int(chapter_number)})

    return data_list


def update_progress(progress, description="", bar_length=20):
    """Display or update a console progress bar.

    Arguments:
    progress    - A float between 0 and 1.  Any int will be converted to a float.
                  A value under 0 represents a "halt." (can be used when
                     an exception has occurred during processing).
                  A value at 1 or bigger represents 100%
    description - The text to be displayed to the left of the status bar.
    bar_length   - How log the bar should be when displayed.
                  Defaults to 20.
    """
    status = ""

    if description == "":
        description = "Status"

    if isinstance(progress, int):
        progress = float(progress)

    if not isinstance(progress, float):
        progress = 0
        status = "error: progress var must be float\r\n"

    if progress < 0:
        progress = 0
        status = "Halt\r\n"

    if progress >= 1:
        progress = 1
        status = "Done\r\n"

    block = int(round(bar_length * progress))
    text = "\r{0}: {1:.2f}% [{2}] {3}".format(description, progress * 100, ("=" * block) + ("-" * (bar_length - block)), status)
    sys.stdout.write(text)
    sys.stdout.flush()


def valid_lecture(lecture_number, lecture_start, lecture_end):
    """Testing if the given lecture number is valid and exist."""
    if lecture_start and lecture_end:
        return lecture_start <= lecture_number <= lecture_end
    elif lecture_start:
        return lecture_start <= lecture_number
    else:
        return lecture_number <= lecture_end


def sanitize_path(path):
    """Cleaning up path for saving files."""
    return "".join([c for c in path if c.isalpha() or c.isdigit() or c in ' .-_,']).rstrip()


def mkdir(directory):
    """Creating output directory structure, if not exist."""
    if not os.path.exists(directory):
        try:
            os.makedirs(directory)
        except OSError as exc:
            if exc.errno != errno.EEXIST:
                raise


def get_caption(caption_file_dir, caption_list):
    """Download captions for Videos."""
    for link, srclang in caption_list:

        dest_dir = os.path.join(".", "captions", caption_file_dir + '-' + srclang.upper() + '.srt')
        mkdir(os.path.dirname(dest_dir))

        if os.path.exists(dest_dir):
            continue
        caption = requests.get(link).text
        caption = re.sub(r"(\d{2,}:\d{2}:\d{2})(\.)(\d{3})", r'\g<1>,\g<3>', caption)
        caption = re.sub(r"WEBVTT\s?\n", "", caption)

        with open(dest_dir, 'w', encoding='utf-8') as save_cap:
            save_cap.write(caption)


def get_attached_file(attached_file_dir, attached_info):
    """Get attached file/supplementary-assets from lecture."""
    course_id = attached_info['course_id']
    lecture_id = attached_info['lecture_id']
    attached_list = attached_info['attached_list']

    for item in attached_list:
        attached_url = ATTACHED_FILE_URL.format(course_id=course_id,
                                                lecture_id=lecture_id,
                                                asset_id=item['id'])

        dest_dir = os.path.join('.', 'attached_files', attached_file_dir, item['filename'])
        mkdir(os.path.dirname(dest_dir))

        if os.path.exists(dest_dir):
            continue

        attached_json = session.get(attached_url).json()
        attached_items = attached_json.get('download_urls')
        
        if debug:
            save_debug_data(attached_json, 'attached_file_' + str(item['id']), 'txt')

        elements = []

        if attached_items is None:
            continue

        for element in attached_items.get('File', {}):
            elements.append(element)

        for element in attached_items.get('SourceCode', {}):
            elements.append(element)

        for element in elements:
            if element.get('label') == 'download' and element.get('file'):
                check_downloaded(element['file'], dest_dir, downloader)


def merge_file(filename, tmp_list):
    """Merge multiple ts file for m3u8."""
    with open(filename, 'wb') as file_mrg:
        for tmp_file in tmp_list:
            with open(tmp_file, 'rb') as merge_tmp:
                shutil.copyfileobj(merge_tmp, file_mrg)
                merge_tmp.close()
            os.remove(tmp_file)


def get_video(link, filename):
    """Send Video link to downloader."""
    m3u8_link = False
    if link.startswith('https://udemy-adaptive-streaming'):
        m3u8_link = True

    if os.path.exists(filename):  # neccessary for m3u8 when ffmpeg not used
        logger.info('%s already exists', filename)
        return

    if use_ffmpeg and m3u8_link:
        check_downloaded(link, filename, 'ffmpeg', do_not_resume=True)

    elif m3u8_link:
        logger.warning('This lecture has m3u8 file.\n'
                       'Please install "ffmpeg" and use argument "--use-ffmpeg"\n'
                       'For better video')
        m3u8_list = []
        m3u8_str = requests.get(link).text

        for line in m3u8_str.splitlines():
            if line.startswith('https://'):
                m3u8_list.append(line)
        total_item = len(m3u8_list)

        for num, url in enumerate(m3u8_list, 1):
            part_name = video_quality + '-' + filename + '.part' + str(num)
            print('\nDownloading {0} of {1} parts:'.format(num, total_item))
            check_downloaded(url, part_name, downloader)

        parts_list = [video_quality + '-' + filename + '.part' + str(num) for num, _ in enumerate(m3u8_list, 1)]
        merge_file(filename, parts_list)
    else:
        tmp_file = video_quality + '-' + filename + '.tmp'
        check_downloaded(link, filename, downloader, tmp_file=tmp_file)


def check_downloaded(link, filename, downloader_name, do_not_resume=False, tmp_file=None):
    """Check if downloading filename already exists."""
    if os.path.exists(filename):
        logger.info('%s already exists', filename)
    else:
        if not tmp_file:
            tmp_file = filename + '.tmp'
        if os.path.exists(tmp_file) and do_not_resume:
            os.remove(tmp_file)
        if 'youtube.com/' in link:
            downloader_name = 'yt_dl'
        download_status = download(link, tmp_file, update_progress, downloader_name)
        if download_status is not 'failed':
            os.rename(tmp_file, filename)


def get_data(directory, filename, link, data_type, attached_info, caption_list):
    """Get lecture content."""
    print('\n\nDownloading {0}: {1!s}\n'.format(data_type, filename))
    previous_dir = os.getcwd()
    mkdir(directory)
    os.chdir(directory)
    file_dir = filename

    if data_type in ['Ebook', 'Audio', 'Presentation', 'File']:
        try:
            delim = link.split('?')[0].rpartition('.')[2]
            if delim in link:
                filename += '.' + delim

            check_downloaded(link, filename, downloader)

        except DLException as exc:
            logger.critical("Couldn't download this %s: %s", data_type, exc)

    elif data_type in ['Video']:
        filename += '.mp4'

        get_video(link, filename)

    elif data_type in ['Pre_pptx']:
        get_pptx(link, filename)

    elif data_type in ['Article']:
        html_json = session.get(link).json()
        html = html_json.get('view_html')
        output_html = '<!DOCTYPE html>\n<html>\n<head>\n<meta charset="UTF-8">\n<title>{0}</title>\n</head>\n<body>\n{1}\n</body>\n</html>'.format(str(filename), html)

        with open(str(filename) + '.html', 'w', encoding='utf-8') as save_art:
            save_art.write(output_html)

    elif data_type in ['quiz']:
        get_quiz(link, filename)

    if attached_info['attached_list']:
        print('\nDownloading attached files\n')
        get_attached_file(file_dir, attached_info)

    if caption_list:
        print('\nDownloading captions\n')
        get_caption(file_dir, caption_list)

    os.chdir(previous_dir)


def get_pptx(img_link, filename):
    """Make pptx from images."""
    pptx_new = pptx.Presentation()
    img_list = [session.get(link) for link in img_link]
    images = [BytesIO(img.content) for img in img_list]
    for img in images:
        slide = pptx_new.slides.add_slide(pptx_new.slide_layouts[6])
        slide.shapes.add_picture(img, 0, 0)

    pptx_new.save("{0}.pptx".format(filename))


def save_link(link, dest, data_type):
    """Save links to a file."""
    if 'Article' in data_type:
        return
    elif 'udemy-adaptive-streaming.udemy.com/' in link:
        logger.warning('Video contains m3u8 file.'
                       'it is better to download instead of saving links')

    with open(dest, 'a', encoding='utf-8') as link_file:
        link_file.write(link + '\n')


def check_course_status(course_id):
    """Check the status of the course."""
    try:
        check_course_status_url = CHECK_COURSE_STATUS_URL.format(course_id=course_id)
        check_course_status_data = session.get(check_course_status_url).text
        if ('user_previewed_course' and 'remaining_seconds') in str(check_course_status_data):
            logger.error('Not a downloadable course.')
            sys.exit(1)
    except TypeError:
        pass


def udemy_dl(username, password, course_link, lecture_start=1, lecture_end=None,
             save_links=False, safe_file_names=False, just_list=False, dest="", dirCourseTitle = False):
    """Login into udemy and do all magic."""
    login(username, password)

    course_id = get_course_id(course_link)

    check_course_status(course_id)

    course_title = get_course_title(course_id)

    if dirCourseTitle:
        parent_directory = '{0!s}'.format(course_title)
        
        if safe_file_names:
            directory = slugify(parent_directory, lower=True, spaces=False, ok='.', only_ascii=True)
        else:
            directory = sanitize_path(parent_directory)
        dest = os.path.join(dest, directory)
        
    last_chapter = -1

    for data in get_data_links(course_id, lecture_start, lecture_end):
        if save_links:
            save_link(data['data_url'], dest, data['data_type'])
        else:
            try:
                directory = '{0:02d} {1!s}'.format(data['chapter_number'], safeencode(data['chapter']))

                if safe_file_names:
                    directory = slugify(directory, lower=True, spaces=False, ok='.', only_ascii=True)
                else:
                    directory = sanitize_path(directory)

            except AttributeError:
                # Fix for untitled opening chapter

                if safe_file_names:
                    directory = '00-opening'
                else:
                    directory = '00 Opening'

            if dest:
                directory = os.path.join(dest, directory)

            filename = '{0:03d} {1!s}'.format(data['lecture_number'], safeencode(data['lecture']))

            if safe_file_names:
                filename = slugify(filename, lower=True, spaces=False, ok='.', only_ascii=True)
            else:
                filename = sanitize_path(filename)

            if just_list:
                if last_chapter != data['chapter_number']:
                    last_chapter = data['chapter_number']
                    print('\r\n{0:02d} {1!s}\r\n=========================='.format(last_chapter, safeencode(data['chapter'])))
                print('{0:03d} {1!s}'.format(data['lecture_number'], safeencode(data['lecture'])))
            else:
                data_url = data['data_url']
                data_type = data['data_type']
                attached_info = data['attached_info']
                caption_list = data['caption_list']
                get_data(directory, filename, data_url, data_type, attached_info, caption_list)

    if os.path.exists(dest) and save_links:
        logger.info('Links successfully saved to : %s\n', os.path.abspath(dest))

    logger.info('Logging out...')
    session.get(LOGOUT_URL)
    logger.info('Logging out success')


def is_integer(num):
    """Check if given value is an integer."""
    try:
        int(num)
        return True
    except ValueError:
        return False


def main():
    """Accepting arguments and preparing."""
    global debug
    global debug_path
    global downloader
    global use_ffmpeg
    global quality_list
    global video_quality

    parser = argparse.ArgumentParser(description='Fetch all the lectures for a udemy course', prog='udemy-dl')
    parser.add_argument('link', help='Link for udemy course', action='store')
    parser.add_argument('-u', '--username', help='Username / Email', default=None, action='store')
    parser.add_argument('-p', '--password', help='Password', default=None, action='store')
    parser.add_argument('--lecture-start', help='Lecture to start at (default is 1)', default=1, action='store')
    parser.add_argument('--lecture-end', help='Lecture to end at (default is last)', default=None, action='store')
    parser.add_argument('-o', '--output', help='Output directory / text file path (if saving links)', default=None, action='store')
    parser.add_argument('-d', '--external-downloader', help='Download with external downloader [aria2c, axel, httpie, curl] (default is aria2c)', default='aria2c', action='store', choices=['aria2c', 'axel', 'httpie', 'curl'])
    parser.add_argument('--use-ffmpeg', help='Download videos from m3u8/hls with ffmpeg (Recommended)', action='store_const', const=True, default=False)
    parser.add_argument('-q', '--video-quality', help='Select video quality [default is 654321(highest)]', default='654321', action='store')
    parser.add_argument('-s', '--save-links', help='Do not download but save links to a file', action='store_const', const=True, default=False)
    parser.add_argument('--safe-file-names', help='Use safe cross-platform filenames', action='store_const', const=True, default=False)
    parser.add_argument('-l', '--list', help='Just list all of the possible lectures and their ids', action='store_const', const=True, default=False)
    parser.add_argument('--debug', help='Enable debug mode', action='store_const', const=True, default=False)
    parser.add_argument('-v', '--version', help='Display the version of udemy-dl and exit', action='version', version='%(prog)s {version}'.format(version=__version__))
    parser.add_argument('--use-course-title', help='Use the course title for the parent folder name (WARNING: can make file path too long)', action='store_const', const=True, default=False)

    args = vars(parser.parse_args())

    username = args['username']
    password = args['password']
    link_args = args['link']
    lecture_start = args['lecture_start']
    lecture_end = args['lecture_end']
    safe_file_names = args['safe_file_names']
    just_list = args['list']
    save_links = args['save_links']
    debug_status = args['debug']
    downloader = args['external_downloader']
    use_ffmpeg = args['use_ffmpeg']
    video_quality = args['video_quality']
    dir_use_course_title = args['use_course_title']

    try:
        if 'udemy.com/draft/' in link_args:
            link = re.search(r'(https://www\.udemy\.com/draft/.*?)/', link_args + '/').group(1)
        else:
            link = re.search(r'(https://www\.udemy\.com/.*?)/', link_args + '/').group(1)
    except AttributeError:
        logger.error('Please follow right url format.\n'
                     'Ex. https://www.udemy.com/COURSE-NAME, https://www.udemy.com/draft/COURSE-ID')
        sys.exit(1)

    course_slug = link.rsplit('/', 1)[1]

    if all(char in '123456' for char in video_quality):
        quality_list = []
        quality_dict = {'1': '240', '2': '360', '3': '480', '4': '640', '5': '720', '6': '1080'}

        for item in video_quality:
            quality_list.append(quality_dict[item])
    else:
        logger.error('Use quality 1 to 6 where 6 being highest.\n'
                     'Ex. 123456(lowest to highest), 3421, 34, 324.\n'
                     'please see usage in wiki\n'
                     'https://github.com/nishad/udemy-dl/wiki/Downloading-lower-video-quality')
        sys.exit(1)

    if use_ffmpeg and not find_executable('ffmpeg'):
        logger.error('ffmpeg is not found. Please install it.')
        sys.exit(1)

    if debug_status:
        debug = True
        debug_path = os.path.abspath(os.path.join(".", 'debug_udemy-dl', time.strftime("%Y%m%d-%H%M%S")))
        mkdir(debug_path)
        logging_name = os.path.join(debug_path, 'debugging.log')
        logging.basicConfig(filename=logging_name, filemode='w', level=0,
                            format='%(asctime)s - [%(levelname)s-%(name)s-%(lineno)d] - %(message)s')
        logger.setLevel(level=logging.DEBUG)
        logger.debug('Debug mode is enabled, debug files will be saved in : \n%s\n', debug_path)
    else:
        debug = False

    if lecture_start is not None:
        if not is_integer(lecture_start) or int(lecture_start) <= 0:
            logger.error('--lecture-start requires natural number argument')
            sys.exit(1)
        lecture_start = int(lecture_start)

    if lecture_end is not None:
        if not is_integer(lecture_end) or int(lecture_end) <= 0:
            logger.error('--lecture-end requires natural number argument')
            sys.exit(1)
        lecture_end = int(lecture_end)

    if args['output']:
        # Normalize the output path if specified
        output_dest = os.path.normpath(args['output'])
        output_file = output_dest
    else:
        # Get output dir name from the URL
        output_dest = os.path.join(".", course_slug)
        output_file = os.path.join(".", course_slug + '.txt')

    if not username:
        username = input("Username / Email : ")

    if not password:
        password = getpass.getpass(prompt='Password : ')

    if save_links:
        # create empty link save location
        if os.path.dirname(output_file) == '':
            output_file = os.path.join(".", output_file)

        mkdir(os.path.dirname(output_file))

        with open(output_file, 'w', encoding='utf-8') as sv_file:
            sv_file.close()
        output_dest = output_file
        logger.info('Saving links to: %s\n', os.path.abspath(output_dest))
    else:
        logger.info('Downloading to: %s\n', os.path.abspath(output_dest))

    udemy_dl(username, password, link, lecture_start, lecture_end,
             save_links, safe_file_names, just_list, output_dest, dir_use_course_title)

if __name__ == '__main__':
    main()
